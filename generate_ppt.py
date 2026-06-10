"""生成答辩 PPT — 选择困难终结 Vibe Coding 答辩"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ═══════════ 配色系统（匹配项目设计规范）═══════════
CREAM      = RGBColor(0xFF, 0xFB, 0xF5)  # 背景
CREAM_DEEP = RGBColor(0xFF, 0xF7, 0xEC)  # 浅卡片
YELLOW_WARM= RGBColor(0xFF, 0xED, 0xB3)  # 暖黄强调
YELLOW_PALE= RGBColor(0xFF, 0xF6, 0xD5)  # 浅黄
PINK_PALE  = RGBColor(0xFF, 0xEF, 0xEC)  # 暖粉
BLACK      = RGBColor(0x1A, 0x1A, 0x1A)  # 黑
SOFT_BLACK = RGBColor(0x3A, 0x3A, 0x3A)  # 软黑
GRAY_MID   = RGBColor(0x7A, 0x7A, 0x7A)  # 灰
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)  # 白
MINT_PALE  = RGBColor(0xED, 0xF8, 0xF0)  # 浅绿
RED        = RGBColor(0xD9, 0x4A, 0x3A)  # 红

FONT = 'Microsoft YaHei'
W = Inches(13.333)  # 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

# ═══════════ 工具函数 ═══════════
def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width or Pt(2.5)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=BLACK):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(3)
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=Pt(14), color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=Pt(12), color=SOFT_BLACK, line_spacing=1.5):
    """lines: list of (text, bold, font_size_override, color_override)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_info in enumerate(lines):
        if isinstance(line_info, str):
            text, b, fs, c = line_info, False, font_size, color
        else:
            text = line_info[0]
            b = line_info[1] if len(line_info) > 1 else False
            fs = line_info[2] if len(line_info) > 2 else font_size
            c = line_info[3] if len(line_info) > 3 else color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = fs
        p.font.color.rgb = c
        p.font.bold = b
        p.font.name = FONT
        p.space_after = Pt(4)
    return txBox

def add_table(slide, left, top, width, height, rows, cols, data, col_widths=None, header_bg=YELLOW_WARM):
    """data: list of lists, first row is header"""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if data[r][c] else ''
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.name = FONT
                p.font.color.rgb = BLACK if r == 0 else SOFT_BLACK
                p.font.bold = (r == 0)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return table_shape

def add_page_number(slide, num):
    add_text_box(slide, Inches(12.2), Inches(7.05), Inches(1), Inches(0.35),
                 str(num), Pt(12), GRAY_MID, False, PP_ALIGN.RIGHT)

def add_card(slide, left, top, width, height, fill=WHITE):
    """添加白底粗黑描边卡片"""
    return add_rect(slide, left, top, width, height, fill, BLACK, Pt(3))

def add_title_bar(slide, title_text, subtitle_text=None):
    """顶部标题条"""
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.15), YELLOW_WARM, BLACK, Pt(3))
    add_text_box(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.65),
                 title_text, Pt(28), BLACK, True)
    if subtitle_text:
        add_text_box(slide, Inches(0.8), Inches(0.72), Inches(11), Inches(0.35),
                     subtitle_text, Pt(13), GRAY_MID, False)

# ═══════════ 第 1 页：封面 ═══════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide, CREAM)

# 顶部装饰条
add_rect(slide, Inches(0), Inches(0), W, Inches(0.15), BLACK)

# 中央内容
add_text_box(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(1.2),
             '🎯 选择困难终结', Pt(48), BLACK, True, PP_ALIGN.CENTER)

# 黄色 pill
pill = add_rounded_rect(slide, Inches(3.8), Inches(2.9), Inches(5.5), Inches(0.7), YELLOW_WARM, BLACK)
add_text_box(slide, Inches(3.8), Inches(2.95), Inches(5.5), Inches(0.6),
             'Vibe Coding · 智能决策助手', Pt(22), BLACK, True, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.0), Inches(11.3), Inches(0.5),
             'Python 程序设计课程大作业  |  全程 AI 生成，零手写代码', Pt(14), GRAY_MID, False, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.8), Inches(11.3), Inches(0.5),
             '线上：dariaaaa.pythonanywhere.com  |  GitHub：dongyifei0502/choice-terminator',
             Pt(12), GRAY_MID, False, PP_ALIGN.CENTER)

# 底部信息
add_text_box(slide, Inches(1), Inches(6.3), Inches(11.3), Inches(0.5),
             '组员：______  ______  ______  |  讲师：李承浩  |  2026-06',
             Pt(14), SOFT_BLACK, False, PP_ALIGN.CENTER)

add_rect(slide, Inches(0), Inches(7.35), W, Inches(0.15), BLACK)
add_page_number(slide, 1)

# ═══════════ 第 2 页：需求设计文档 ═══════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, CREAM)
add_title_bar(slide, '需求设计文档', '答辩要求 1 — 需求背景与功能设计')

# 左侧：需求背景
add_card(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(2.4))
add_multiline(slide, Inches(0.9), Inches(1.55), Inches(5), Inches(2.1), [
    ('📌 需求背景', True, Pt(16), BLACK),
    '',
    '• 日常选择困难症：纠结午餐/购物/旅行目的地',
    '• 需要一个多维度加权打分系统，用数据辅助决策',
    '• 需要云端存储历史记录，支持多设备同步',
    '• 用户无需注册即可使用核心决策功能',
], Pt(12.5))

# 右侧：设计文档引用
add_card(slide, Inches(6.6), Inches(1.4), Inches(6.2), Inches(2.4))
add_multiline(slide, Inches(7.0), Inches(1.55), Inches(5.5), Inches(2.1), [
    ('📐 设计文档（全部 AI 生成）', True, Pt(16), BLACK),
    '',
    ('📄 需求梳理.md — 项目背景、目标用户、核心功能', False, Pt(12)),
    ('📄 PRD-选择困难终结.md — 数据模型、算法、验收标准', False, Pt(12)),
    ('📄 开发计划.md — 12步依赖关系 + 完成标准checklist', False, Pt(12)),
    ('📄 UI设计规范.md — CSS变量表、组件矩阵、响应式断点', False, Pt(12)),
], Pt(12))

# 底部：功能需求表格
add_card(slide, Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.8))
add_text_box(slide, Inches(0.9), Inches(4.25), Inches(3), Inches(0.4), '功能模块总览', Pt(16), BLACK, True)

data = [
    ['模块', '核心功能', '优先级'],
    ['决策引擎', '输入选项 → 场景选择 → 多轮打分 → 加权计算 → 推荐结果', 'P0'],
    ['结果展示', '推荐Pill + 柱状图(Chart.js) + 理由(规则引擎) + 得分明细表', 'P0'],
    ['用户系统', '注册/登录/JWT认证(零外部依赖:SHA-256+HMAC)', 'P2'],
    ['历史记录', '云端SQLite同步 + localStorage兜底 + JSON/CSV导出', 'P1'],
    ['管理后台', '问题库CRUD(系统+自定义) + 用户封禁管理', 'P2'],
]
add_table(slide, Inches(0.9), Inches(4.7), Inches(11.5), Inches(2.0),
          len(data), 3, data, [Inches(1.8), Inches(7.6), Inches(1.0)])

add_page_number(slide, 2)

# ═══════════ 第 3 页：系统架构 ═══════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, CREAM)
add_title_bar(slide, '系统架构 — 逻辑先行', 'Vibe Coding 要求 1 — 先画架构，再让 AI 生成')

# 技术架构
add_card(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(2.8))
add_text_box(slide, Inches(0.9), Inches(1.55), Inches(5), Inches(0.4), '🏗️ 技术架构', Pt(16), BLACK, True)
# 架构框图
frontend_box = add_rounded_rect(slide, Inches(1.2), Inches(2.15), Inches(4.5), Inches(1.05), WHITE, BLACK)
add_multiline(slide, Inches(1.5), Inches(2.25), Inches(4), Inches(0.85), [
    ('前端 SPA (Vanilla JS)', True, Pt(13), BLACK),
    ('首页→问答→评分→结果→历史→管理', False, Pt(11)),
    ('Chart.js / CSS3 / HTML5', False, Pt(10.5), GRAY_MID),
])
# 箭头
add_text_box(slide, Inches(3.2), Inches(3.25), Inches(0.8), Inches(0.4), 'REST API', Pt(10), GRAY_MID, False, PP_ALIGN.CENTER)
backend_box = add_rounded_rect(slide, Inches(1.2), Inches(3.5), Inches(4.5), Inches(0.6), WHITE, BLACK)
add_text_box(slide, Inches(1.4), Inches(3.55), Inches(4.1), Inches(0.5),
             'Python Flask · SQLite', Pt(12), BLACK, True, PP_ALIGN.CENTER)

# 右侧：依赖关系
add_card(slide, Inches(6.6), Inches(1.4), Inches(6.2), Inches(2.8))
add_text_box(slide, Inches(7.0), Inches(1.55), Inches(5), Inches(0.4), '🔗 12步依赖关系', Pt(16), BLACK, True)

# ASCII 图
add_multiline(slide, Inches(7.0), Inches(2.15), Inches(5.5), Inches(1.9), [
    ('S0脚手架→S1首页→S2问答页→S3评分引擎', False, Pt(10.5), SOFT_BLACK),
    ('                    ↓                      ', False, Pt(10.5), SOFT_BLACK),
    ('              S4结果页←S5图表+S6理由        ', False, Pt(10.5), SOFT_BLACK),
    ('                    ↓                      ', False, Pt(10.5), SOFT_BLACK),
    ('        S7历史→S8导出  S9管理→S10打磨       ', False, Pt(10.5), SOFT_BLACK),
    ('                    ↓                      ', False, Pt(10.5), SOFT_BLACK),
    ('          阶段三Flask后端→阶段四封禁系统      ', False, Pt(10.5), SOFT_BLACK),
    ('                    ↓                      ', False, Pt(10.5), SOFT_BLACK),
    ('              阶段五PythonAnywhere部署       ', False, Pt(10.5), BLACK),
], Pt(10.5))

# 底部：Vibe Coding 三步
add_card(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5))
add_text_box(slide, Inches(0.9), Inches(4.65), Inches(5), Inches(0.4),
             '🎯 Vibe Coding 三步原则', Pt(16), BLACK, True)

steps = [
    ('① 逻辑先行', '写任一行代码前，先用\n文字+流程图梳理核心逻辑', CREAM_DEEP),
    ('② 渐进式', '分模块分步骤生成，\n不一次写整个系统', YELLOW_PALE),
    ('③ 责任归属', '小组对运行代码负100%\n责任，AI写错不是理由', PINK_PALE),
]
for i, (title, desc, bg) in enumerate(steps):
    x = Inches(0.9 + i * 3.9)
    add_card(slide, x, Inches(5.15), Inches(3.5), Inches(1.6), bg)
    add_text_box(slide, x + Inches(0.3), Inches(5.25), Inches(3.0), Inches(0.4), title, Pt(15), BLACK, True)
    add_text_box(slide, x + Inches(0.3), Inches(5.7), Inches(3.0), Inches(0.9), desc, Pt(11.5), SOFT_BLACK)

add_page_number(slide, 3)

# ═══════════ 第 4 页：小组分工 ═══════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, CREAM)
add_title_bar(slide, '小组分工', '三人协作 · Vibe Coding 角色分配')

add_card(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(2.0))
add_text_box(slide, Inches(0.9), Inches(1.55), Inches(5), Inches(0.4), '角色分配', Pt(16), BLACK, True)
data = [
    ['组员', '角色', '负责模块', 'AI 工具'],
    ['组员 A', '需求+架构师', 'PRD撰写、逻辑梳理、12步开发计划、提示词设计、分步分解', 'Claude Code'],
    ['组员 B', '后端+数据', 'Flask骨架、认证系统、SQLite数据库、API接口、权限中间件', 'Claude Code'],
    ['组员 C', '前端+部署', 'SPA页面开发、CSS样式系统、Chart.js图表、PythonAnywhere部署', 'Claude Code'],
]
add_table(slide, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.4),
          len(data), 4, data, [Inches(1.3), Inches(1.5), Inches(6.5), Inches(1.5)])

add_card(slide, Inches(0.5), Inches(3.7), Inches(12.3), Inches(2.8))
add_text_box(slide, Inches(0.9), Inches(3.85), Inches(5), Inches(0.4), '各阶段参与度', Pt(16), BLACK, True)
data2 = [
    ['阶段', 'A（架构）', 'B（后端）', 'C（前端）'],
    ['需求分析', '★ 主导', '参与讨论', '参与讨论'],
    ['12步开发计划', '★ 主导', '提供反馈', '提供反馈'],
    ['Flask后端', '审核API设计', '★ 主导', '—'],
    ['前端6页面', '—', '—', '★ 主导'],
    ['四维改进', '★ 方案设计', '配合调整', '配合调整'],
    ['用户+封禁', '—', '★ 主导', '前端对接'],
    ['PythonAnywhere部署', '—', '协助排查', '★ 主导'],
]
add_table(slide, Inches(0.9), Inches(4.3), Inches(11.5), Inches(2.0),
          len(data2), 4, data2, [Inches(3.0), Inches(3.0), Inches(3.0), Inches(2.5)])

add_page_number(slide, 4)

# ═══════════ 第 5 页：系统演示 ═══════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, CREAM)
add_title_bar(slide, '系统演示', '答辩要求 2 — 6个演示点 · 约5分钟')

add_card(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(3.2))
add_text_box(slide, Inches(0.9), Inches(1.55), Inches(5), Inches(0.4), '演示路径', Pt(16), BLACK, True)
data = [
    ['顺序', '演示内容', '时长', '展示重点'],
    ['1', '打开网站 → 首页场景选择 + 输入选项', '30s', '4种场景切换、标签解析'],
    ['2', '问答页：5道题逐题打分 + 进度条', '1min', '滑块交互、权重Badge、前后导航'],
    ['3', '结果页：推荐 + 柱状图 + 理由 + 明细', '1min', '暖黄推荐pill、Chart.js动画'],
    ['4', '平局场景演示：🎲 随机 + 换一个', '30s', '同分随机选、换一个按钮'],
    ['5', '登录(admin/admin123) → 管理问题库', '1min', 'CRUD操作、场景筛选'],
    ['6', '用户管理：封禁/解封', '1min', '管理员专属、实时生效'],
]
add_table(slide, Inches(0.9), Inches(2.0), Inches(11.5), Inches(2.5),
          len(data), 4, data, [Inches(0.6), Inches(4.5), Inches(0.9), Inches(5.0)])

add_card(slide, Inches(0.5), Inches(4.9), Inches(12.3), Inches(2.0))
add_text_box(slide, Inches(0.9), Inches(5.05), Inches(5), Inches(0.4), '演示要点', Pt(16), BLACK, True)
add_multiline(slide, Inches(0.9), Inches(5.5), Inches(11.5), Inches(1.2), [
    '🎯 展示 4 种决策场景（餐饮/购物/旅行/通用），切换场景时 placeholder 联动变化',
    '👤 展示用户注册 → 自动登录 → 云端历史同步流程',
    '⚙️ 展示管理页 Tab 切换（问题管理 / 用户管理），来源列区分"系统"与"自定义"',
    '🔒 展示管理员专属权限：系统题不可被普通用户修改删除',
], Pt(12.5))

add_page_number(slide, 5)

# ═══════════ 第 6 页：最棒的 Prompt ═══════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, CREAM)
add_title_bar(slide, '最棒的 Prompt', '答辩要求 3-展示1 — 精确引导 AI 写出正确代码')

prompt_text = (
    '"你现在是 Step 2：问答页的开发。依赖 Step 0 和 Step 1。请实现：'
    '1. 创建 js/router.js — 页面切换函数（display 显隐 + init 函数）\n'
    '2. 创建 js/services/questionService.js — selectQuestions() 按权重排序选题\n'
    '3. 创建 js/components/progressBar.js — updateProgress(current, total)\n'
    '4. 创建 js/components/slider.js — 滑块渲染（默认值3）+\n'
    '   bindSliderEvents() 实时更新分数（>=4 加 .high 类）+ collectScores()\n'
    '5. 修改 index.html 为多页面容器\n'
    '6. 修改 app.js 整合路由 + 问答页逻辑\n'
    '完成标准：首页输入3选项→进入问答→进度条+权重badge→滑块打分→前后导航→最后一题变粉色"'
)

add_card(slide, Inches(0.5), Inches(1.4), Inches(7.8), Inches(5.5))
add_text_box(slide, Inches(0.9), Inches(1.55), Inches(5), Inches(0.4), 'Prompt 原文', Pt(16), BLACK, True)
add_multiline(slide, Inches(0.9), Inches(2.1), Inches(7.0), Inches(4.6), [
    (prompt_text, False, Pt(11), SOFT_BLACK),
], Pt(11))

# 右侧：为什么有效
add_card(slide, Inches(8.6), Inches(1.4), Inches(4.2), Inches(5.5))
add_text_box(slide, Inches(9.0), Inches(1.55), Inches(4), Inches(0.4), '为什么有效', Pt(16), BLACK, True)

reasons = [
    ('① 明确前置依赖', '"依赖Step 0和Step 1" → AI不会重复造轮子'),
    ('② 分模块编号', '6个子任务清晰独立 → AI不会遗漏'),
    ('③ 指定函数签名', 'updateProgress(current,total) → 入参类型被约束'),
    ('④ 附完成标准', '验收条件可自我验证 → AI输出更精准'),
    ('⑤ 上下文完整', '单次会话独立理解 → 不需要之前对话'),
]
y = 2.1
for title, desc in reasons:
    add_text_box(slide, Inches(9.0), Inches(y), Inches(3.5), Inches(0.3), title, Pt(12), YELLOW_WARM, True)
    add_text_box(slide, Inches(9.0), Inches(y + 0.3), Inches(3.5), Inches(0.5), desc, Pt(10.5), SOFT_BLACK)
    y += 0.9

add_page_number(slide, 6)

# ═══════════ 第 7 页：AI 犯错案例 ═══════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, CREAM)
add_title_bar(slide, 'AI 犯错案例 ×3', '答辩要求 3-展示2 — 发现AI幻觉并引导修复')

# 案例 1
add_card(slide, Inches(0.5), Inches(1.4), Inches(3.9), Inches(5.5))
add_text_box(slide, Inches(0.7), Inches(1.55), Inches(3.5), Inches(0.4), '❌ 案例1：循环导入', Pt(14), BLACK, True)
add_multiline(slide, Inches(0.7), Inches(2.0), Inches(3.5), Inches(4.7), [
    ('AI 写的代码：', True, Pt(11)),
    ('# app.py', False, Pt(10), GRAY_MID),
    ('from backend.auth import auth_bp', False, Pt(10), RED),
    ('', False, Pt(8)),
    ('# backend/auth.py', False, Pt(10), GRAY_MID),
    ('from app import get_db  ← 循环!', False, Pt(10), RED),
    ('', False, Pt(8)),
    ('发现：', True, Pt(11)),
    ('启动报 ImportError →\nPython循环导入经典错误', False, Pt(10.5)),
    ('', False, Pt(8)),
    ('修复：', True, Pt(11)),
    ('提取 db.py 独立模块,\napp.py 和 auth.py\n都从 backend.db 导入', False, Pt(10.5)),
], Pt(10.5))

# 案例 2
add_card(slide, Inches(4.7), Inches(1.4), Inches(3.9), Inches(5.5))
add_text_box(slide, Inches(4.9), Inches(1.55), Inches(3.5), Inches(0.4), '❌ 案例2：参数顺序', Pt(14), BLACK, True)
add_multiline(slide, Inches(4.9), Inches(2.0), Inches(3.5), Inches(4.7), [
    ('AI 写的代码：', True, Pt(11)),
    ('Store.setUser(', False, Pt(10)),
    ('  data.user.username,', False, Pt(10), RED),
    ('  data.user.role,', False, Pt(10), RED),
    ('  data.token', False, Pt(10), RED),
    (');', False, Pt(10)),
    ('', False, Pt(8)),
    ('实际签名：', True, Pt(11)),
    ('setUser(token, username, role)', False, Pt(10.5)),
    ('                     ↑ 顺序反了!', False, Pt(10.5), RED),
    ('', False, Pt(8)),
    ('发现：', True, Pt(11)),
    ('登录后用户名显示为\ntoken 字符串 → 排查到\n参数传反', False, Pt(10.5)),
    ('', False, Pt(8)),
    ('修复：', True, Pt(11)),
    ('全局替换参数顺序', False, Pt(10.5)),
], Pt(10.5))

# 案例 3
add_card(slide, Inches(8.9), Inches(1.4), Inches(3.9), Inches(5.5))
add_text_box(slide, Inches(9.1), Inches(1.55), Inches(3.5), Inches(0.4), '❌ 案例3：依赖缺失', Pt(14), BLACK, True)
add_multiline(slide, Inches(9.1), Inches(2.0), Inches(3.5), Inches(4.7), [
    ('AI 使用了：', True, Pt(11)),
    ('flask-cors, bcrypt, pyjwt', False, Pt(11), RED),
    ('', False, Pt(8)),
    ('PythonAnywhere 免费版', False, Pt(10.5)),
    ('无法 pip install 这些包', False, Pt(10.5), RED),
    ('', False, Pt(8)),
    ('发现：', True, Pt(11)),
    ('部署后 HTTP 500,\n错误日志空白 →\n逐个排查依赖', False, Pt(10.5)),
    ('', False, Pt(8)),
    ('修复 Prompt：', True, Pt(11)),
    ('"用纯标准库替代：', False, Pt(10.5)),
    ('hashlib→bcrypt', False, Pt(10.5)),
    ('hmac→PyJWT', False, Pt(10.5)),
    ('手写after_request→CORS"', False, Pt(10.5)),
    ('最终 requirements.txt\n只需要 flask', False, Pt(10.5)),
], Pt(10.5))

add_page_number(slide, 7)

# ═══════════ 第 8 页：Vibe Coding 使用反思 ═══════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, CREAM)
add_title_bar(slide, 'Vibe Coding 使用反思', '答辩要求 4 — 好的方面、不足、核心体会')

# 好的方面
add_card(slide, Inches(0.5), Inches(1.4), Inches(6.0), Inches(2.8))
add_text_box(slide, Inches(0.9), Inches(1.55), Inches(5), Inches(0.4), '✅ 好的方面', Pt(16), BLACK, True)
data_good = [
    ['优势', '说明'],
    ['🚀 效率极高', '12步前端+4步后端+部署，1天内完成，传统至少1周'],
    ['📐 结构清晰', 'AI严格按照开发计划的依赖关系执行，不跳过依赖'],
    ['🔧 改起来快', '"四维改进"四个问题一次对话全部解决'],
    ['🌐 部署全程AI', '从render.yaml到PythonAnywhere API上传文件'],
]
add_table(slide, Inches(0.9), Inches(2.0), Inches(5.2), Inches(2.0),
          len(data_good), 2, data_good, [Inches(1.2), Inches(4.0)])

# 需要改进
add_card(slide, Inches(6.8), Inches(1.4), Inches(6.0), Inches(2.8))
add_text_box(slide, Inches(7.2), Inches(1.55), Inches(5), Inches(0.4), '⚠️ 需要改进', Pt(16), BLACK, True)
data_bad = [
    ['不足', '说明'],
    ['⚠️ AI幻觉', 'AI会用不存在的API(如PythonAnywhere console API)'],
    ['🧪 测试依赖人', 'AI无法自动打开浏览器测试，验证全靠人工'],
    ['📝 Prompt质量', '描述不清时AI写出"看起来对但实际不完整"的代码'],
    ['🔐 安全需审查', '认证/加密代码需人工确认逻辑正确性'],
]
add_table(slide, Inches(7.2), Inches(2.0), Inches(5.2), Inches(2.0),
          len(data_bad), 2, data_bad, [Inches(1.2), Inches(4.0)])

# 核心体会
add_card(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5), CREAM_DEEP)
add_text_box(slide, Inches(0.9), Inches(4.65), Inches(5), Inches(0.4), '💡 核心体会', Pt(16), BLACK, True)

add_multiline(slide, Inches(0.9), Inches(5.15), Inches(11.5), Inches(1.6), [
    ('Vibe Coding ≠ 不思考。恰恰相反，它要求更强的「架构思维」和「Prompt 设计能力」。', True, Pt(15), BLACK),
    ('', False, Pt(6)),
    ('你不需要写代码，但你必须：', False, Pt(13), SOFT_BLACK),
    ('  1. 把大问题拆成小步骤（分模块、分步骤）→ 开发计划.md 就是12个AI能消化的粒度', False, Pt(12)),
    ('  2. 每一步清晰定义输入、输出、边界条件 → 每步的"目标+涉及文件+完成标准"格式', False, Pt(12)),
    ('  3. 亲自验证每一步的输出是否正确 → 启动→登录→打分→部署，全流程人工测试', False, Pt(12)),
    ('  4. 对最终代码负 100% 责任 → "AI写错了"不是报错理由，你必须找到错误并修好它', False, Pt(12)),
], Pt(12))

add_page_number(slide, 8)

# ═══════════ 第 9 页：Vibe Coding 完整规划过程 ═══════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, CREAM)
add_title_bar(slide, 'Vibe Coding 完整规划过程', '从零到一的5步规划 — 全部由AI辅助完成')

# 5步流程表格
add_card(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(2.0))
add_text_box(slide, Inches(0.9), Inches(1.55), Inches(5), Inches(0.4), '5步规划流程', Pt(16), BLACK, True)
data_plan = [
    ['顺序', '产出物', 'AI做了什么', '我做了什么'],
    ['①', '需求梳理.md', '根据场景描述生成结构化需求文档', '描述场景、目标用户、核心功能'],
    ['②', 'PRD-选择困难终结.md', '扩展为完整PRD(数据模型/算法/验收标准)', '审核功能定义、补充边界条件'],
    ['③', '开发计划.md', '拆解为12步依赖关系图+完成标准', '确认依赖关系、调整步长'],
    ['④', 'demo.html', '生成静态设计Demo(5页面+Chart.js)', '检查视觉风格、调整配色布局'],
    ['⑤', 'UI设计规范.md', '从Demo提取完整设计规范文档', '确认规范一致性'],
]
add_table(slide, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.4),
          len(data_plan), 4, data_plan, [Inches(0.5), Inches(2.5), Inches(4.5), Inches(4.0)])

# ③ 开发计划：关键一步
add_card(slide, Inches(0.5), Inches(3.7), Inches(6.0), Inches(3.3), YELLOW_PALE)
add_text_box(slide, Inches(0.9), Inches(3.85), Inches(5), Inches(0.4),
             '⭐ 最关键的一步：开发计划', Pt(15), BLACK, True)
add_multiline(slide, Inches(0.9), Inches(4.3), Inches(5.3), Inches(2.5), [
    ('这是把 PRD 变成 AI 可执行指令的核心转化：', False, Pt(12), SOFT_BLACK),
    ('', False, Pt(4)),
    ('① 拆解粒度：3500行代码分成12步，每步AI能一次消化', False, Pt(11.5)),
    ('② 完成标准：每步末尾附checklist，变成自动验收条件', False, Pt(11.5)),
    ('③ 文件标注：[新建]/[修改] 让AI知道操作类型', False, Pt(11.5)),
    ('④ 独立可执行：每个步骤描述足够完整，新AI会话可直接实施', False, Pt(11.5)),
    ('', False, Pt(4)),
    ('"每个步骤要足够完整，让新的AI会话能独立理解并执行"', True, Pt(12), BLACK),
], Pt(11.5))

# 规划→编码
add_card(slide, Inches(6.8), Inches(3.7), Inches(6.0), Inches(3.3))
add_text_box(slide, Inches(7.2), Inches(3.85), Inches(5), Inches(0.4),
             '🚀 规划 → 编码：一句命令', Pt(15), BLACK, True)
add_multiline(slide, Inches(7.2), Inches(4.4), Inches(5.3), Inches(2.2), [
    ('规划全部完成后...', False, Pt(13), GRAY_MID),
    ('', False, Pt(8)),
    ('"从1开始，逐步推进"', True, Pt(22), BLACK),
    ('', False, Pt(8)),
    ('AI 就按照开发计划.md 的依赖关系顺序，', False, Pt(12), SOFT_BLACK),
    ('一步一验收，逐步实现了整个项目。', False, Pt(12), SOFT_BLACK),
    ('', False, Pt(8)),
    ('约50轮 Prompt，1天完成从零到上线。', True, Pt(13), BLACK),
], Pt(12))

add_page_number(slide, 9)

# ═══════════ 第 10 页：优势、不足、自评分 ═══════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, CREAM)
add_title_bar(slide, '优势、不足、自评分', '答辩必答题')

# 优势
add_card(slide, Inches(0.5), Inches(1.4), Inches(6.0), Inches(2.8))
add_text_box(slide, Inches(0.9), Inches(1.55), Inches(5), Inches(0.4), '🏆 我们组的优势', Pt(16), BLACK, True)
add_multiline(slide, Inches(0.9), Inches(2.1), Inches(5.3), Inches(2.0), [
    ('🎯 需求明确：PRD+UI规范+12步计划，每步有完成标准', False, Pt(12)),
    ('🏗️ 架构完整：前端SPA+Flask后端+SQLite+云端部署', False, Pt(12)),
    ('🔄 渐进式：严格按依赖关系分步推进→基础→功能→增强→部署', False, Pt(12)),
    ('💬 Prompt规范：每步带"目标+文件+做什么+完成标准"四要素', False, Pt(12)),
    ('🐛 问题记录完善：循环导入/参数错误/依赖缺失等7个典型问题', False, Pt(12)),
], Pt(12))

# 不足
add_card(slide, Inches(6.8), Inches(1.4), Inches(6.0), Inches(2.8))
add_text_box(slide, Inches(7.2), Inches(1.55), Inches(5), Inches(0.4), '🔧 我们的不足', Pt(16), BLACK, True)
add_multiline(slide, Inches(7.2), Inches(2.1), Inches(5.3), Inches(2.0), [
    ('📱 无移动端适配：仅700px响应式，未做手机端专项适配', False, Pt(12)),
    ('🧪 无自动化测试：完全依赖人工验证，无单元测试覆盖', False, Pt(12)),
    ('🔒 安全简化：bcrypt→SHA-256，PyJWT→HMAC，企业级不够', False, Pt(12)),
    ('🛒 无邮箱/找回密码：用户系统缺少账户恢复机制', False, Pt(12)),
    ('📊 未接AI推荐理由：Step 11(AI增强)未实施，仅规则模板', False, Pt(12)),
], Pt(12))

# 自评分
add_card(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5), YELLOW_PALE)
add_text_box(slide, Inches(0.9), Inches(4.65), Inches(5), Inches(0.4), '📊 自评分', Pt(16), BLACK, True)
add_multiline(slide, Inches(0.9), Inches(5.15), Inches(11.5), Inches(1.6), [
    ('评级：A 档', True, Pt(28), BLACK),
    ('', False, Pt(6)),
    ('理由：', True, Pt(14), BLACK),
    ('• 系统完整可用：6页面SPA + 14个API端点 + 云端部署 + 公网可访问', False, Pt(13)),
    ('• 遵循Vibe Coding三步要求：逻辑先行(5份规划文档)→渐进式(12步逐层构建)→责任归属(7个AI错误被找出并修复)', False, Pt(13)),
    ('• 有详细的Prompt记录(第6页)和AI纠错案例(第7页)，小组分工明确(第4页)', False, Pt(13)),
], Pt(13))

add_page_number(slide, 10)

# ═══════════ 第 11 页：答辩分工 ═══════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, CREAM)
add_title_bar(slide, '答辩分工', '每人承担一个环节 · 合计 5-8 分钟')

add_card(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(3.0))
add_text_box(slide, Inches(0.9), Inches(1.55), Inches(5), Inches(0.4), '分工表', Pt(16), BLACK, True)
data_div = [
    ['环节', '负责人', '页码', '内容', '时长'],
    ['需求/设计文档', '组员 A', '第2-3页', '需求背景→架构设计→12步依赖关系', '~1.5min'],
    ['Vibe Coding规划', '组员 A', '第9页', '5步规划过程(需求→PRD→计划→Demo→规范)', '~1min'],
    ['系统演示', '组员 B', '第5页', '6个演示点全程操作展示', '~2min'],
    ['Prompt+纠错', '组员 C', '第6-7页', '最棒Prompt分析+3个AI犯错案例', '~1.5min'],
    ['反思+必答题', '组员 A', '第8/10页', '反思+优势/不足+自评A档理由', '~1min'],
]
add_table(slide, Inches(0.9), Inches(2.0), Inches(11.5), Inches(2.2),
          len(data_div), 5, data_div, [Inches(1.8), Inches(1.2), Inches(1.0), Inches(5.0), Inches(1.0)])

add_card(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.3))
add_text_box(slide, Inches(0.9), Inches(4.85), Inches(5), Inches(0.4), '关键数据', Pt(16), BLACK, True)
data_stats = [
    ['指标', '数值', '指标', '数值', '指标', '数值'],
    ['项目文件', '31个', '总代码行数', '~3500行', '前端JS', '15个文件'],
    ['后端Python', '6个文件', '数据库表', '3张', 'API端点', '14个'],
    ['Git提交', '13次', 'Prompt轮次', '~50轮', '外部依赖', '仅Flask+Chart.js'],
]
add_table(slide, Inches(0.9), Inches(5.3), Inches(11.5), Inches(1.5),
          len(data_stats), 6, data_stats,
          [Inches(1.5), Inches(1.8), Inches(1.5), Inches(1.8), Inches(1.5), Inches(2.0)])

add_page_number(slide, 11)

# ═══════════ 保存 ═══════════
output_path = r'c:\Users\Daria\Desktop\choice terminator\答辩PPT-选择困难终结.pptx'
prs.save(output_path)
print('PPT saved: ' + output_path)
print('Slides: ' + str(len(prs.slides)))